import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
from PIL import Image

class SyntaxHighlightingText(ScrolledText):
    def __init__(self, master=None, **kwargs):
        ScrolledText.__init__(self, master, **kwargs)
        self._highlight_pattern_id = None
        self._configure_tags()

    def _configure_tags(self):
        self.tag_configure("NewSlide", foreground="blue")
        self.tag_configure("FormatNumber", foreground="darkred")
        self.tag_configure("Title", foreground="forest green")
        self.tag_configure("Content", foreground="darkorchid")
        self.tag_configure("Image", foreground="chocolate")
        self.tag_configure("ImagePosition", foreground="chocolate")
        self.tag_configure("ImageSize", foreground="chocolate")

    def highlight_pattern(self, pattern, tag, start="1.0", end="end", regexp=False):
        start = self.index(start)
        end = self.index(end)
        self.mark_set("matchStart", start)
        self.mark_set("matchEnd", start)
        self.mark_set("searchLimit", end)
        count = tk.IntVar()

        while True:
            index = self.search(pattern, "matchEnd", "searchLimit", count=count, regexp=regexp)
            if index == "":
                break
            self.mark_set("matchStart", index)
            self.mark_set("matchEnd", f"{index}+{count.get()}c")
            self.tag_add(tag, "matchStart", "matchEnd")

    def highlight_syntax(self):
        keywords = {
            "NewSlide": "NewSlide",
            "FormatNumber": "FormatNumber",
            "Title": "Title",
            "Content": "Content",
            "Image": "Image",
            "ImagePosition": "ImagePosition",
            "ImageSize": "ImageSize"
        }

        for tag, pattern in keywords.items():
            self.highlight_pattern(pattern, tag, regexp=True)
            self.highlight_pattern(pattern.lower(), tag, regexp=True)
            self.highlight_pattern(pattern.replace('_', ''), tag, regexp=True)
            self.highlight_pattern(pattern.replace('_', '').lower(), tag, regexp=True)

        self.after(100, self.highlight_syntax)

class SlideCompileEditor:
    def __init__(self, root):
        self.root = root
        self.root.title("SlideCompile")
        
        self.file_path = None

        self.text_area = SyntaxHighlightingText(self.root, wrap='word')
        self.text_area.pack(fill='both', expand=True)
        
        self.menu_bar = tk.Menu(self.root)
        self.file_menu = tk.Menu(self.menu_bar, tearoff=0)
        
        self.file_menu.add_command(label="Open", command=self.open_file)
        self.file_menu.add_command(label="Save", command=self.save_file)
        self.file_menu.add_command(label="Save As", command=self.save_as_file)
        self.file_menu.add_separator()
        self.file_menu.add_command(label="Exit", command=self.root.quit)
        
        self.menu_bar.add_cascade(label="File", menu=self.file_menu)
        self.root.config(menu=self.menu_bar)
        
        self.root.bind_all('<Control-s>', self.handle_save_shortcut)
        self.root.bind_all('<Control-c>', self.handle_copy_shortcut)

        self.compile_button = tk.Button(self.root, text="Compile", command=self.compile)
        self.compile_button.place(relx=1.0, rely=1.0, anchor='se', bordermode='outside')

        self.status_label = tk.Label(self.root, text="", fg="blue")
        self.status_label.pack(side='bottom')

        self.text_area.highlight_syntax()

    def handle_save_shortcut(self, event):
        self.save_file()

    def handle_copy_shortcut(self, event):
        selected_text = self.text_area.get(tk.SEL_FIRST, tk.SEL_LAST)
        self.root.clipboard_clear()
        self.root.clipboard_append(selected_text)
        
    def open_file(self):
        try:
            self.file_path = filedialog.askopenfilename(defaultextension=".txt", 
                                               filetypes=[("Text files", "*.txt"), ("All files", "*.*")])
            if self.file_path:
                with open(self.file_path, "r") as file:
                    content = file.read()
                    self.text_area.delete(1.0, tk.END)
                    self.text_area.insert(tk.INSERT, content)
                    self.root.title(f"SlideCompile - {self.file_path}")
        except FileNotFoundError:
            self.show_error("File not found. Please select a valid file.")
        except Exception as e:
            self.show_error(f"Could not open file: {e}")

    def save_file(self):
        if self.file_path:
            try:
                with open(self.file_path, "w") as file:
                    content = self.text_area.get(1.0, tk.END)
                    file.write(content)
                    self.root.title(f"SlideCompile - {self.file_path}")
            except Exception as e:
                self.show_error(f"Could not save file: {e}")
        else:
            self.save_as_file()

    def save_as_file(self):
        try:
            self.file_path = filedialog.asksaveasfilename(defaultextension=".txt", 
                                                          filetypes=[("Text files", "*.txt"), ("All files", "*.*")])
            if self.file_path:
                with open(self.file_path, "w") as file:
                    content = self.text_area.get(1.0, tk.END)
                    file.write(content)
                    self.root.title(f"SlideCompile - {self.file_path}")
        except Exception as e:
            self.show_error(f"Could not save file: {e}")

    def add_slide(self, prs, format_number, title, content, image_url, image_position, image_size, line):
        try:
            slide_layout = prs.slide_layouts[format_number-1]
            slide = prs.slides.add_slide(slide_layout)

            if title:
                slide.shapes.title.text = title

            if content:
                bullet_points = content.split('\n')
                content_box = slide.shapes.placeholders[1]
                content_frame = content_box.text_frame
                for point in bullet_points:
                    p = content_frame.add_paragraph()
                    p.text = point

            if image_url:
                response = requests.get(image_url)
                if response.status_code == 200:
                    image_data = BytesIO(response.content)
                    img = Image.open(image_data)
                    image_path = f"temp_image_{line}.png"
                    img.save(image_path)

                    width, height = self.get_image_size(image_size)
                    left, top = self.get_dynamic_image_position(image_position, width, height)
                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                else:
                    raise ImageDownloadError(image_url, line)

        except IndexError:
            raise SlideFormatError(format_number, line)
        except requests.RequestException:
            raise ImageDownloadError(image_url, line)
        except Exception as e:
            raise SlideCreationError(line, str(e))

    def get_dynamic_image_position(self, position, width, height):
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        
        positions = {
            'top left': (Inches(0.5), Inches(0.5)),
            'top right': (slide_width - width - Inches(0.5), Inches(0.5)),
            'bottom left': (Inches(0.5), slide_height - height - Inches(0.5)),
            'bottom right': (slide_width - width - Inches(0.5), slide_height - height - Inches(0.5)),
            'middle left': (Inches(0.5), (slide_height - height) / 2),
            'middle right': (slide_width - width - Inches(0.5), (slide_height - height) / 2),
            'middle': ((slide_width - width) / 2, (slide_height - height) / 2)
        }
        return positions.get(position.lower(), (Inches(1), Inches(1)))

    def get_image_size(self, size):
        sizes = {
            'tiny': (Inches(1), Inches(1)),
            'small': (Inches(2), Inches(2)),
            'medium': (Inches(4), Inches(3)),
            'large': (Inches(6), Inches(4.5)),
            'extra large': (Inches(10), Inches(7.5))
        }
        return sizes.get(size.lower(), (Inches(4), Inches(3)))

    def compile(self):
        self.save_file()
        self.status_label.config(text="Compiling the slideshow, this may take a while if you have lots of images...")
        self.root.update_idletasks()

        try:
            with open(self.file_path) as file:
                lines = file.readlines()
        except FileNotFoundError:
            self.show_error("File not found. Please save the file first.")
            return

        prs = Presentation()

        try:
            self.current_slide_data = {'format_number': -1, 'title': '', 'content': '', 'image_url': '', 'image_position': '', 'image_size': 'medium'}
            for lineNumber, line in enumerate(lines, start=1):
                line = line.strip()

                if line.lower() == "newslide" or line.lower() == "new_slide":
                    if self.current_slide_data['format_number'] != -1:
                        self.add_slide(prs, self.current_slide_data['format_number'], self.current_slide_data['title'], self.current_slide_data['content'], self.current_slide_data['image_url'], self.current_slide_data['image_position'], self.current_slide_data['image_size'], lineNumber)
                    self.current_slide_data = {'format_number': -1, 'title': '', 'content': '', 'image_url': '', 'image_position': '', 'image_size': 'medium'}
                elif line.lower().startswith("formatnumber:") or line.lower().startswith("format_number"):
                    _, format_number = line.split(':', 1)
                    self.current_slide_data['format_number'] = int(format_number.strip())
                elif line.lower().startswith("title:"):
                    _, title = line.split(':', 1)
                    self.current_slide_data['title'] = title.strip()
                elif line.lower().startswith("content:"):
                    _, content = line.split(':', 1)
                    self.current_slide_data['content'] += content.strip() + '\n'
                elif line.lower().startswith("image:"):
                    _, image_url = line.split(':', 1)
                    self.current_slide_data['image_url'] = image_url.strip()
                elif line.lower().startswith("imageposition:") or line.lower().startswith("image_position:"):
                    _, image_position = line.split(':', 1)
                    self.current_slide_data['image_position'] = image_position.strip()
                elif line.lower().startswith("imagesize:") or line.lower().startswith("image_size:"):
                    _, image_size = line.split(':', 1)
                    self.current_slide_data['image_size'] = image_size.strip()

            if self.current_slide_data['format_number'] != -1:
                self.add_slide(prs, self.current_slide_data['format_number'], self.current_slide_data['title'], self.current_slide_data['content'], self.current_slide_data['image_url'], self.current_slide_data['image_position'], self.current_slide_data['image_size'], lineNumber)
        except Exception as e:
            self.show_error(f"An unexpected error occurred: {e}")
            self.status_label.config(text="")
            return

        file_path = filedialog.asksaveasfilename(defaultextension=".pptx")
        if file_path:
            prs.save(file_path)

        self.status_label.config(text="Compilation finished!")
        
    def show_error(self, error_message):
        messagebox.showerror("Error", error_message)
        self.status_label.config(text="")

class SlideCompileError(Exception):
    pass


class SlideFormatError(SlideCompileError):
    def __init__(self, format_number, line_number):
        self.format_number = format_number
        self.line_number = line_number
        self.message = f"Invalid slide format number '{format_number}' on line {line_number}."
        super().__init__(self.message)

class ImageDownloadError(SlideCompileError):
    def __init__(self, image_url, line_number):
        self.image_url = image_url
        self.line_number = line_number
        self.message = f"Could not download image from '{image_url}' on line {line_number}. Please check the image URL."
        super().__init__(self.message)

class SlideCreationError(SlideCompileError):
    def __init__(self, line_number, error):
        self.line_number = line_number
        self.error = error
        self.message = f"Error creating slide on line {line_number}: {error}. Please check the slide data."
        super().__init__(self.message)

if __name__ == "__main__":
    root = tk.Tk()
    app = SlideCompileEditor(root)
    root.mainloop()